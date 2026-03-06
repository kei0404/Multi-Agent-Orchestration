#!/usr/bin/env python3
"""
PowerPoint presentation generator with company branding.
- Logo: upper-right of every slide
- Footer: ⓈSYSTEMS DESIGN Co.Ltd copyright
"""

import os
import sys
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE


# --- Branding constants ---
FOOTER_TEXT = "\u24c8systems DESIGN Co.Ltd"
LOGO_WIDTH = Inches(1.2)
LOGO_MARGIN_RIGHT = Inches(0.3)
LOGO_MARGIN_TOP = Inches(0.2)
FOOTER_FONT_SIZE = Pt(8)
FOOTER_FONT_COLOR = RGBColor(0x66, 0x66, 0x66)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def find_logo():
    """Search for logo file in known locations."""
    assets_dir = Path(__file__).parent.parent / "assets"
    candidates = [
        assets_dir / "sdcロゴ（透明）.png",
        assets_dir / "logo.png",
        assets_dir / "logo.jpg",
        assets_dir / "logo.svg",
        Path.cwd() / "logo.png",
        Path.cwd() / "logo.jpg",
    ]
    for p in candidates:
        if p.exists():
            return str(p)
    return None


def add_logo(slide, prs, logo_path):
    """Add company logo to upper-right corner of the slide."""
    if not logo_path or not os.path.exists(logo_path):
        return

    slide_w = prs.slide_width
    logo_w = LOGO_WIDTH
    left = slide_w - logo_w - LOGO_MARGIN_RIGHT
    top = LOGO_MARGIN_TOP

    slide.shapes.add_picture(logo_path, left, top, width=logo_w)


def add_footer(slide, prs):
    """Add copyright footer to the bottom of the slide."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    footer_h = Inches(0.4)
    left = Inches(0.5)
    top = slide_h - footer_h - Inches(0.1)
    width = slide_w - Inches(1.0)

    txBox = slide.shapes.add_textbox(left, top, width, footer_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    year = datetime.now().year
    run.text = f"\u00a9 {year} {FOOTER_TEXT}"
    run.font.size = FOOTER_FONT_SIZE
    run.font.color.rgb = FOOTER_FONT_COLOR


def add_branding(slide, prs, logo_path):
    """Apply both logo and footer branding to a slide."""
    add_logo(slide, prs, logo_path)
    add_footer(slide, prs)


def create_title_slide(prs, title, subtitle=None, logo_path=None):
    """Create a title slide (first slide of presentation)."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle

    add_branding(slide, prs, logo_path)
    return slide


def create_content_slide(prs, title, body_lines, logo_path=None):
    """Create a slide with title and bullet points."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line

    add_branding(slide, prs, logo_path)
    return slide


def create_section_slide(prs, title, logo_path=None):
    """Create a section header slide."""
    layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    add_branding(slide, prs, logo_path)
    return slide


def create_blank_slide(prs, logo_path=None):
    """Create a blank slide with branding only."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    add_branding(slide, prs, logo_path)
    return slide


def create_two_column_slide(prs, title, left_lines, right_lines, logo_path=None):
    """Create a two-column content slide."""
    layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    # Left column
    left_body = slide.placeholders[1]
    tf_left = left_body.text_frame
    tf_left.clear()
    for i, line in enumerate(left_lines):
        if i == 0:
            tf_left.paragraphs[0].text = line
        else:
            tf_left.add_paragraph().text = line

    # Right column
    right_body = slide.placeholders[2]
    tf_right = right_body.text_frame
    tf_right.clear()
    for i, line in enumerate(right_lines):
        if i == 0:
            tf_right.paragraphs[0].text = line
        else:
            tf_right.add_paragraph().text = line

    add_branding(slide, prs, logo_path)
    return slide


def build_presentation(slide_data, output_path, logo_path=None):
    """
    Build a full presentation from structured slide data.

    slide_data: list of dicts, each with:
        - type: "title" | "content" | "section" | "two_column" | "blank"
        - title: str
        - subtitle: str (title slide only)
        - body: list[str] (content slide)
        - left: list[str], right: list[str] (two_column)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for sd in slide_data:
        stype = sd.get("type", "content")

        if stype == "title":
            create_title_slide(prs, sd["title"], sd.get("subtitle"), logo_path)
        elif stype == "content":
            create_content_slide(prs, sd["title"], sd.get("body", []), logo_path)
        elif stype == "section":
            create_section_slide(prs, sd["title"], logo_path)
        elif stype == "two_column":
            create_two_column_slide(
                prs, sd["title"],
                sd.get("left", []), sd.get("right", []),
                logo_path
            )
        elif stype == "blank":
            create_blank_slide(prs, logo_path)

    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path


# --- CLI entry point ---
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create_pptx.py <slides.json> [output.pptx]")
        print()
        print("slides.json format:")
        print(json.dumps([
            {"type": "title", "title": "Presentation Title", "subtitle": "Subtitle"},
            {"type": "content", "title": "Slide Title", "body": ["Point 1", "Point 2"]},
        ], indent=2, ensure_ascii=False))
        sys.exit(1)

    json_path = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    with open(json_path, "r", encoding="utf-8") as f:
        slides = json.load(f)

    logo = find_logo()
    build_presentation(slides, output, logo)
